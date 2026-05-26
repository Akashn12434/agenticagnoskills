#!/usr/bin/env python3
"""Create a PowerPoint presentation from a topic and slide outline."""

from __future__ import annotations

import argparse
import json
import os
import platform
import re
import subprocess
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_DIR = BASE_DIR / "output"

#Color palette
NAVY = RGBColor(28, 40, 51)
SLATE = RGBColor(46, 64, 83)
TEAL = RGBColor(17, 122, 101)
LIGHT_GRAY = RGBColor(213, 216, 220)
WHITE = RGBColor(255, 255, 255)

#Clean the text, swap all spaces and punctuation for underscores, trim the edges so it can be used as a filename
def slugify(value: str) -> str:
    """Return a filesystem-friendly slug."""
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", value.strip().lower()).strip("_")
    return slug or "presentation"

#Automatically opens the generated PowerPoint.
def open_file(path: Path) -> None:
    """Open the generated presentation with the system default application."""
    if platform.system() == "Windows":
        os.startfile(path)  # type: ignore[attr-defined]
    elif platform.system() == "Darwin":
        subprocess.Popen(["open", str(path)])
    else:
        subprocess.Popen(["xdg-open", str(path)])

#This creates fallback slides, If AI does not supply slides:
def default_slides(topic: str) -> list[dict[str, Any]]:
    return [
        {
            "title": topic,
            "subtitle": "An overview presentation",
            "bullets": [],
        },
        {
            "title": "Introduction",
            "bullets": [
                f"What {topic} means and why it matters",
                "Key terms and background context",
                "Main ideas covered in this presentation",
            ],
        },
        {
            "title": "Topic Overview",
            "bullets": [
                "Important components and relationships",
                "How the topic is commonly organized",
                "Essential facts to remember",
            ],
        },
        {
            "title": "Key Concepts",
            "bullets": [
                "Core concept one with a short explanation",
                "Core concept two with a short explanation",
                "Core concept three with a short explanation",
            ],
        },
        {
            "title": "Examples",
            "bullets": [
                "Representative example or case",
                "How the example connects to the main idea",
                "What learners should notice",
            ],
        },
        {
            "title": "Practical Applications",
            "bullets": [
                "Where the topic appears in real life",
                "Current uses or future opportunities",
                "Why understanding it is useful",
            ],
        },
        {
            "title": "Summary",
            "bullets": [
                "Recap of the most important points",
                "Final takeaway for the audience",
                "Questions and discussion",
            ],
        },
    ]


def unwrap_slide(slide: dict[str, Any]) -> dict[str, Any]:
    """Accept common nested model outputs and return the actual slide dict."""
    if "title" in slide or "bullets" in slide or "subtitle" in slide:
        return slide

    for key in ("slide", "content", "data", "example_key"):
        nested = slide.get(key)
        if isinstance(nested, dict):
            return unwrap_slide(nested)

    if len(slide) == 1:
        nested = next(iter(slide.values()))
        if isinstance(nested, dict):
            return unwrap_slide(nested)

    return slide

#cleans slide data. It ensures that each slide has a title, subtitle, and bullets, and applies defaults if any of these are missing. It also limits the number of bullets to 6 and trims whitespace.
def normalize_slides(topic: str, slides: list[dict[str, Any]] | None) -> list[dict[str, Any]]:
    if not slides:
        return default_slides(topic)

    normalized: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        slide = unwrap_slide(slide)
        title = str(slide.get("title") or f"Slide {index}").strip() 
        subtitle = str(slide.get("subtitle") or "").strip()
        raw_bullets = slide.get("bullets") or []
        bullets = [str(item).strip() for item in raw_bullets if str(item).strip()]
        normalized.append({"title": title, "subtitle": subtitle, "bullets": bullets[:6]})

    if normalized:
        normalized[0]["subtitle"] = normalized[0].get("subtitle") or "An overview presentation"
    return normalized

#Reusable font formatter.Instead of repeating:everywhere, we can just call this function to apply consistent styling to all text elements in the presentation.
def set_text(run, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold

#Creates slide numbers. Bottom-right:
def add_footer(slide, slide_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.78), Inches(12.2), Inches(0.35))
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = f"{slide_number}"
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.runs[0].font.size = Pt(15)
    paragraph.runs[0].font.color.rgb = SLATE

#Creates Slide 1. This is the cover page.
def add_title_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY

    #Thin teal line across top.
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    #Adds title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.05), Inches(11.8), Inches(1.35))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    paragraph = title_frame.paragraphs[0]
    paragraph.text = slide_data["title"]
    set_text(paragraph.runs[0], 42, WHITE, True)

    #An overview presentation subtitle, default if not provided by AI.
    subtitle = slide_data.get("subtitle") or "An overview presentation"
    subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(3.48), Inches(9.8), Inches(0.55))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    set_text(subtitle_frame.paragraphs[0].runs[0], 22, LIGHT_GRAY)

    #Adds the date to the bottom-left of the slide.
    date_box = slide.shapes.add_textbox(Inches(0.78), Inches(6.35), Inches(5.0), Inches(0.35))
    date_frame = date_box.text_frame
    date_frame.text = date.today().isoformat()
    set_text(date_frame.paragraphs[0].runs[0], 12, LIGHT_GRAY)

    add_footer(slide, 1)

#Creates every slide after the title slide.
def add_content_slide(prs: Presentation, slide_data: dict[str, Any], slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(11.9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = slide_data["title"]
    set_text(title_frame.paragraphs[0].runs[0], 32, NAVY, True)

    underline = slide.shapes.add_shape(1, Inches(0.62), Inches(1.28), Inches(2.1), Inches(0.07))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()

    content_box = slide.shapes.add_textbox(Inches(0.88), Inches(1.72), Inches(11.35), Inches(4.75))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    bullets = slide_data.get("bullets") or ["Key point"]
    for index, bullet in enumerate(bullets[:6]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"★ {bullet}"
        paragraph.level = 0
        paragraph.space_after = Pt(9)
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(22 if len(bullets) <= 4 else 20)
        paragraph.font.color.rgb = SLATE

    add_footer(slide, slide_number)


def create_presentation(
    topic: str,
    slides: list[dict[str, Any]] | None = None,
    output_path: str | None = None,
    auto_open: bool = True,
) -> str:
    """Create a .pptx file and return its absolute path."""
    #1. Normalize and clean the slide data to ensure it has a consistent structure and defaults.
    normalized_slides = normalize_slides(topic, slides)
    #2. Create a new PowerPoint presentation and set the slide dimensions.
    prs = Presentation()
    #3. Set widescreen format (16:9 aspect ratio) which is standard for modern presentations.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    #4. Add the title slide using the first slide's data, then add content slides for the rest of the slides.
    add_title_slide(prs, normalized_slides[0])
    for slide_number, slide_data in enumerate(normalized_slides[1:], start=2):
        add_content_slide(prs, slide_data, slide_number)
    #5. Save the presentation to the specified output path 
    if output_path:
        output = Path(output_path)
        if not output.is_absolute():
            output = OUTPUT_DIR / output
    else:
        output = OUTPUT_DIR / f"{slugify(topic)}.pptx"

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    resolved_output = output.resolve()
    if auto_open:
        open_file(resolved_output)
    return str(resolved_output)

#Helper function to load slides from a JSON string or file path. This allows users to provide slide data in a flexible way, either directly as a JSON string or as a path to a JSON file.
def load_slides(value: str | None) -> list[dict[str, Any]] | None:
    if not value:
        return None
    path = Path(value)
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return json.loads(value)

#The main function parses command-line arguments for the presentation topic, optional slides data, output path, and whether to auto-open the presentation. 
def main() -> None:
    parser = argparse.ArgumentParser(description="Create a PowerPoint deck.")
    parser.add_argument("topic", help="Presentation topic")
    parser.add_argument("--slides", help="JSON slide list or path to a JSON file")
    parser.add_argument("--output", help="Output .pptx path")
    parser.add_argument("--no-open", action="store_true", help="Do not open the presentation after creating it")
    args = parser.parse_args()

    result = create_presentation(args.topic, load_slides(args.slides), args.output, auto_open=not args.no_open)
    print(result)


if __name__ == "__main__":
    main()
